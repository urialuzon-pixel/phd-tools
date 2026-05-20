from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

DARK_BLUE = RGBColor(0x1A, 0x37, 0x6C)
ACCENT    = RGBColor(0x2E, 0x86, 0xC1)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF4, 0xF6, 0xF9)
GRAY      = RGBColor(0x55, 0x6B, 0x82)

W, H = Inches(13.33), Inches(7.5)


def set_rtl(tf):
    for para in tf.paragraphs:
        pPr = para._pPr
        if pPr is None:
            pPr = para._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
        para.alignment = PP_ALIGN.RIGHT


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=WHITE,
             rtl=True, align=PP_ALIGN.RIGHT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Arial"
    if rtl:
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=DARK_BLUE, indent=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        pPr.set(qn('a:rtl'), '1')
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Arial"
    return txBox


prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, DARK_BLUE)
add_rect(s, 0, Inches(2.8), W, Inches(2.0), ACCENT)

add_text(s, "הצעת מחקר לדוקטורט",
         Inches(1), Inches(0.4), Inches(11), Inches(0.8),
         font_size=14, color=RGBColor(0xAA,0xCC,0xFF))

add_text(s, "אי-ענישה גיאופוליטית ופעילות כרייה",
         Inches(0.5), Inches(1.1), Inches(12), Inches(1.0),
         font_size=32, bold=True, color=WHITE)

add_text(s, "מסגרת היברידית Physics–ML לגילוי אנומליות כרייה בשטחים שנויים במחלוקת",
         Inches(0.5), Inches(2.0), Inches(12), Inches(0.7),
         font_size=20, color=RGBColor(0xCC,0xDD,0xFF))

add_text(s, "הגשה פנימית — טיוטה ראשונה",
         Inches(1), Inches(3.1), Inches(11), Inches(0.5),
         font_size=14, color=WHITE)

add_text(s, "אוריאל אוזון  |  2026",
         Inches(1), Inches(6.5), Inches(11), Inches(0.5),
         font_size=13, color=RGBColor(0xAA,0xCC,0xFF))

# ── Slide 2: The Problem ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "הבעיה", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

bullets = [
    "כאשר שלטון ריבוני שנוי במחלוקת — מנגנוני הפיקוח על הכרייה קורסים",
    "קונפליקט מוביל לאי-אוכלוסייה: תושבים עוזבים, כוחות אכיפה נסוגים",
    "נותרת קרקע עשירה במשאבים עם מינימום התנגדות לכניסה",
    "Lambin & Meyfroidt (2011): לחץ שלטוני מזיז פעילות כרייה לאזורים בעלי פיקוח חלש",
    "שטחים שנויים במחלוקת = קצה הסקאלה של ריקנות שלטונית",
    "Prem et al. (2020): קריסת FARC בקולומביה → גל כרייה בלתי מפוקחת",
    "הפער: אין מסגרת גילוי שיטתית, סקלאבילית, ובזמן אמת",
]
add_bullets(s, bullets, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5),
            font_size=17, color=DARK_BLUE)

# ── Slide 3: Research Questions ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "שאלות המחקר", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

add_rect(s, Inches(0.4), Inches(1.3), Inches(12.4), Inches(1.5), ACCENT)
add_text(s, "RQ0 (ראשית): האם תנאי אי-ענישה גיאופוליטית — קרבה לסכסוכי ריבונות, גירעונות שלטוניים, עוצמת קונפליקט — מנבאים שאריות עכירות חריגות (לוויין מינוס פיזיקה) באגני נהרות שנויים במחלוקת?",
         Inches(0.5), Inches(1.35), Inches(12), Inches(1.4),
         font_size=14, bold=False, color=WHITE)

rqs = [
    "RQ1: האם מסגרת Physics-ML היברידית (SWAT + XGBoost/LSTM) מגלה אנומליות כרייה באגנים עם מחסור בנתונים תוך שליטה בשונות הידרולוגית טבעית?",
    "RQ2: האם משתנים גיאופוליטיים (MID/ICOW, מדדי שלטון, ACLED) מוסיפים כוח ניבוי מובהק מעבר למנבאים סביבתיים בלבד?",
    "RQ3: האם מודל משחקי רב-שחקנים מכוייל מתבניות השאריות מעתיק את עיתוי ועוצמת אנומליות הכרייה ההיסטוריות?",
]
add_bullets(s, rqs, Inches(0.5), Inches(3.0), Inches(12), Inches(4.0),
            font_size=15, color=DARK_BLUE)

# ── Slide 4: Methodological Innovation ───────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "החידוש המתודולוגי", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

add_text(s, "גישות קיימות", Inches(7.5), Inches(1.3), Inches(5.3), Inches(0.5),
         font_size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER, rtl=False)
add_text(s, "הגישה הזו", Inches(0.5), Inches(1.3), Inches(5.3), Inches(0.5),
         font_size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, rtl=False)

existing = [
    "× סיווג כיסוי קרקע / זיהוי שינוי",
    "× קו בסיס שרירותי",
    "× ללא הקשר גיאופוליטי",
    "× ללא זיהוי סיבתי",
    "× מטריצת תמורות מונחה הנחות",
]
add_bullets(s, existing, Inches(7.3), Inches(1.9), Inches(5.5), Inches(3.5),
            font_size=15, color=GRAY)

new_approach = [
    "✓ שארית מכוילת פיזיקלית כאות גילוי",
    "✓ SWAT כנגד-עובדה מדוד",
    "✓ משתנים שלטוניים וקונפליקט כמנבאים",
    "✓ DiD מדורג לזיהוי סיבתי",
    "✓ מטריצת תמורות מהאות הלוויין עצמו",
]
add_bullets(s, new_approach, Inches(0.5), Inches(1.9), Inches(6.5), Inches(3.5),
            font_size=15, color=DARK_BLUE)

add_rect(s, Inches(0.4), Inches(5.6), Inches(12.4), Inches(1.4),
         RGBColor(0xE8, 0xF4, 0xFD))
add_text(s,
         "ההבחנה המרכזית: במקום לסווג כרייה, אנו מודדים חריגה מקו בסיס פיזיקלי — "
         "מה שהופך את האנומליה לפרשנית ביחידות פיזיקליות ועמידה בפני בלבולים סביבתיים",
         Inches(0.6), Inches(5.7), Inches(12), Inches(1.1),
         font_size=14, color=DARK_BLUE)

# ── Slide 5: 6-Step Pipeline ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "הצינור האנליטי: 6 שלבים", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

steps = [
    ("1", "קו בסיס פיזיקלי", "מודל SWAT → צפי SSC יומי"),
    ("2", "תצפית לוויינית", "Sentinel-2 NDTI → עכירות נצפית"),
    ("3", "חישוב שאריות", "נצפה מינוס צפוי, מנורמל"),
    ("4", "מנבא ML", "XGBoost + LSTM: סביבתי בלבד vs. + גיאופוליטי"),
    ("5", "ולידציה סיבתית", "DiD מדורג: הסלמה כטיפול"),
    ("6", "מודל תיאורטי", "משחק רב-שחקנים מכוייל מהנתונים"),
]

box_w, box_h = Inches(1.8), Inches(0.95)
gap = Inches(0.22)
start_x = Inches(0.35)
start_y = Inches(1.45)

for i, (num, title, desc) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    add_rect(s, x, start_y, box_w, box_h, DARK_BLUE)
    add_text(s, num, x, start_y + Inches(0.05), box_w, Inches(0.35),
             font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=False)
    add_text(s, title, x, start_y + Inches(0.4), box_w, Inches(0.5),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)

    add_rect(s, x, start_y + box_h + Inches(0.08),
             box_w, Inches(1.0), RGBColor(0xE8,0xF4,0xFD))
    add_text(s, desc, x + Inches(0.05),
             start_y + box_h + Inches(0.1),
             box_w - Inches(0.1), Inches(0.9),
             font_size=12, color=DARK_BLUE,
             align=PP_ALIGN.CENTER, rtl=True)

add_rect(s, Inches(0.4), Inches(4.6), Inches(12.4), Inches(0.5), ACCENT)
add_text(s, "שלבים 1–3: אות הגילוי  |  שלבים 4–5: מבנה הסברתי וסיבתי  |  שלב 6: מודל מנגנוני",
         Inches(0.5), Inches(4.65), Inches(12), Inches(0.4),
         font_size=13, color=WHITE, align=PP_ALIGN.CENTER, rtl=True)

add_text(s,
         "עיקרון מנחה: השארית הפיזיקלית (שלב 3) היא האות המרכזי — כל השאר בונה עליה",
         Inches(0.5), Inches(5.3), Inches(12), Inches(0.6),
         font_size=15, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, rtl=True)

add_text(s,
         "יתרון מרכזי: ללא תלות בנתוני מד-מים פנימיים — חיוני לאגנים שנויים במחלוקת",
         Inches(0.5), Inches(6.1), Inches(12), Inches(0.5),
         font_size=14, color=GRAY, align=PP_ALIGN.CENTER, rtl=True)

# ── Slide 6: Three Communities ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "ייחוד התרומה — שלוש קהילות מחקר", Inches(0.5), Inches(0.2),
         Inches(12), Inches(0.8), font_size=28, bold=True, color=WHITE)

cols = [
    ("חישה מרחוק\nוגילוי כרייה", DARK_BLUE,
     ["✓ גילוי ראשון המשתמש\nבשארית מכוילת פיזיקלית",
      "✓ לא רק סיווג — אנומליה\nביחידות פיזיקליות",
      "✓ חסין לבלבולים\nסביבתיים"]),
    ("ML פיזיקלי\nלהידרולוגיה", ACCENT,
     ["✓ ראשון להכניס\nמשתני שלטון/קונפליקט\nכמנבאים לשאריות",
      "✓ SHAP מזהה תרומה\nמוסף של ההקשר הפוליטי",
      "✓ מחיל DiD מדורג\nעל תוצאה רציפה"]),
    ("כלכלה פוליטית\nשל כרייה", RGBColor(0x17,0x6D,0x3D),
     ["✓ ראשון לכייל מטריצת\nתמורות משחקית\nמאות לוויין פיזיקלי",
      "✓ מגשר בין\nחישה מרחוק ו-IR",
      "✓ מאפשר תרחישי\nסיכון עתידיים"]),
]

col_w = Inches(3.9)
for j, (title, col, items) in enumerate(cols):
    x = Inches(0.4) + j * (col_w + Inches(0.25))
    add_rect(s, x, Inches(1.3), col_w, Inches(0.7), col)
    add_text(s, title, x, Inches(1.32), col_w, Inches(0.65),
             font_size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, rtl=True)
    add_rect(s, x, Inches(2.05), col_w, Inches(4.9), RGBColor(0xF0,0xF4,0xFF))
    add_bullets(s, items, x + Inches(0.1), Inches(2.15),
                col_w - Inches(0.2), Inches(4.6),
                font_size=14, color=DARK_BLUE)

# ── Slide 7: Data Sources ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "מקורות נתונים", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

data_cols = [
    ("לוויין", [
        "Sentinel-2 MSI (ESA)",
        "10 מטר, 5 ימים",
        "NDTI דרך Google Earth Engine",
    ]),
    ("הידרולוגיה", [
        "SRTM DEM",
        "FAO soil + Copernicus",
        "ERA5 climate reanalysis",
        "ללא תלות במד-מים",
    ]),
    ("גיאופוליטיקה", [
        "MID — סכסוכים בינמדינתיים",
        "ICOW — תביעות טריטוריאליות",
        "ACLED — אירועי קונפליקט",
        "World Governance Indicators",
    ]),
    ("גיאולוגיה", [
        "BGS World Mineral Statistics",
        "USGS mineral assessments",
        "כמשתנה כלי לאנדוגניות",
    ]),
]

col_w = Inches(3.0)
for j, (title, items) in enumerate(data_cols):
    x = Inches(0.35) + j * (col_w + Inches(0.28))
    add_rect(s, x, Inches(1.3), col_w, Inches(0.55), ACCENT)
    add_text(s, title, x, Inches(1.33), col_w, Inches(0.5),
             font_size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, rtl=True)
    add_rect(s, x, Inches(1.9), col_w, Inches(4.6),
             RGBColor(0xF0,0xF4,0xFF))
    add_bullets(s, items, x + Inches(0.1), Inches(2.0),
                col_w - Inches(0.2), Inches(4.3),
                font_size=14, color=DARK_BLUE)

# ── Slide 8: Case Selection ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "בחירת מקרי הבוחן", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

criteria = [
    "קריטריונים לבחירה:",
    "  • מיקסום שונות בסוג הסכסוך, משטר השלטון, עושר מינרלי",
    "  • עמידה בסף Sentinel-2: כיסוי עננים < 40% מהחודשים",
    "  • זמינות נתוני ACLED/MID מקבילים (עם תאריכי הסלמה ספציפיים)",
    "",
    "אזורים מועמדים:",
    "  • אגן האמזונס — סכסוכי גבול, כרייה שלא ברישיון",
    "  • מסדרון הסאהל — ריבונות חלשה, משאבים מינרליים",
    "  • דרום-מזרח אסיה — שטחים שנויים במחלוקת, פעילות כרייה",
    "",
    "בחירה סופית: בהמשך לביקורת זמינות נתונים — שנה ראשונה",
]
add_bullets(s, criteria, Inches(0.5), Inches(1.4), Inches(12), Inches(5.5),
            font_size=16, color=DARK_BLUE)

# ── Slide 9: Timeline ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, LIGHT_BG)
add_rect(s, 0, 0, W, Inches(1.2), DARK_BLUE)
add_text(s, "לוח זמנים", Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=28, bold=True, color=WHITE)

timeline = [
    ("שנה 1 — ס׳ א׳", "בחירת מקרים + ביקורת נתונים; הגדרת SWAT ל-2-3 אגנים; צינור GEE"),
    ("שנה 1 — ס׳ ב׳", "חישוב NDTI וסדרות זמן שאריות (שלבים 2-3); מפות אנומליה ראשוניות"),
    ("שנה 2 — ס׳ א׳", "מנבא ML (שלב 4): XGBoost + LSTM; שילוב משתנים גיאופוליטיים; SHAP"),
    ("שנה 2 — ס׳ ב׳", "ולידציה סיבתית (שלב 5): DiD מדורג; הרחבה לכלל האגנים"),
    ("שנה 3 — ס׳ א׳", "מודל תיאורטי (שלב 6): ספציפיקציה, כיול, ניתוח שיווי משקל"),
    ("שנה 3 — ס׳ ב׳", "אינטגרציה; עיבוד; פרסום מדיניות; שחרור קוד פתוח"),
]

row_h = Inches(0.82)
for i, (period, activity) in enumerate(timeline):
    y = Inches(1.35) + i * row_h
    bg_col = RGBColor(0xE8,0xF4,0xFD) if i % 2 == 0 else WHITE
    add_rect(s, Inches(0.3), y, Inches(12.7), row_h - Inches(0.05), bg_col)
    add_text(s, period, Inches(9.0), y + Inches(0.15), Inches(3.8), Inches(0.6),
             font_size=13, bold=True, color=ACCENT, align=PP_ALIGN.RIGHT, rtl=True)
    add_text(s, activity, Inches(0.5), y + Inches(0.15), Inches(8.2), Inches(0.6),
             font_size=13, color=DARK_BLUE, align=PP_ALIGN.RIGHT, rtl=True)

# ── Slide 10: Open Questions ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_bg(s, DARK_BLUE)
add_rect(s, 0, 0, W, Inches(1.2), ACCENT)
add_text(s, "שאלות פתוחות וסיכונים — לדיון עם המנחה",
         Inches(0.5), Inches(0.2), Inches(12), Inches(0.8),
         font_size=24, bold=True, color=WHITE)

questions = [
    "? זיהוי סיבתי: האם קיימת מספיק שונות בתאריכי ההסלמה לחוזק DiD מדורג?",
    "? סיבתיות הפוכה: האם כרייה גוררת קונפליקט, ולא רק להיפך?",
    "? תוקף חיצוני: האם 3 אזורים מספיקים לכלליות?",
    "? מדידה גיאופוליטית: ACLED/MID לוכד את הריקנות השלטונית הרלוונטית?",
    "? סיכון שלב 6: האם מודל משחקים תלוי בהצלחת שלבים 1-5 — מה הסיכון המצטבר?",
    "? הגדרת אנומליה: מהו הסף המינימלי לכרייה 'אמיתית' לעומת רעש טבעי?",
]
add_bullets(s, questions, Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.5),
            font_size=16, color=WHITE)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\urial\OneDrive\Dev\phd_research_agent\presentations\PhD_Proposal_Hebrew.pptx"
prs.save(out)
print(f"Saved: {out}")
